import os
import sys
import io
import zipfile
import tempfile
import shutil
import time
import json
import requests
from datetime import datetime
import streamlit as st
import pandas as pd

# Path setup to ensure imports work correctly
ROOT = os.path.dirname(os.path.abspath(__file__))
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)

# ── Playwright auto-install on Streamlit Cloud ────────────────────────────────
def _ensure_playwright():
    """Tries to install Playwright Chromium on first run (Cloud-safe)."""
    try:
        import subprocess
        result = subprocess.run(
            ["playwright", "install", "chromium", "--with-deps"],
            capture_output=True, text=True, timeout=120
        )
    except Exception:
        pass

# ── Import backend modules ────────────────────────────────────────────────────
from backend.github.analyzer import analyze, AnalyzerError
import backend.github.github_client as github_client_mod
import backend.deployment.vercel_client as vercel_client_mod
import backend.ai.ai_engine as ai_engine
import backend.ai.content_writer as content_writer_mod
from backend.ai.project_summary import build_summary, ProjectSummaryError
from backend.screenshot.capture import capture_screenshots

# ── Page Config ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="DevLaunch AI — Automation Agent",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── Dynamic Reinitialization Helper ───────────────────────────────────────────
def reinit_clients(groq_key: str, github_token: str, vercel_token: str = None):
    """Binds UI key inputs directly to the backend client singletons."""
    if groq_key:
        os.environ["GROQ_API_KEY"] = groq_key
        try:
            ai_engine.engine = ai_engine.AIEngine(api_key=groq_key)
            content_writer_mod.writer = content_writer_mod.ContentWriter()
            import backend.screenshot.formatter as formatter_mod
            formatter_mod.formatter = formatter_mod.Formatter()
        except Exception as e:
            st.sidebar.error(f"AI Engine Error: {e}")

    if github_token:
        os.environ["GITHUB_TOKEN"] = github_token
        try:
            github_client_mod.github_client = github_client_mod.GitHubClient(token=github_token)
        except Exception as e:
            st.sidebar.error(f"GitHub Client Error: {e}")

    if vercel_token:
        os.environ["VERCEL_TOKEN"] = vercel_token
        try:
            vercel_client_mod.vercel_client = vercel_client_mod.VercelClient(token=vercel_token)
        except Exception:
            pass

# ── Git Cloning Helper ────────────────────────────────────────────────────────
def clone_repo(repo_url: str, dest_dir: str) -> bool:
    """Clones a remote repository to a local destination directory."""
    import subprocess
    try:
        result = subprocess.run(
            f"git clone {repo_url} {dest_dir}",
            shell=True, capture_output=True, text=True, timeout=120
        )
        return result.returncode == 0
    except Exception:
        return False

# ── URL Metadata Scraper Helper ───────────────────────────────────────────────
def scrape_url_metadata(url: str) -> dict:
    """Scrapes meta tags, headings, and text content from a URL for AI analysis."""
    try:
        headers = {"User-Agent": "Mozilla/5.0 (compatible; DevLaunchAI/1.0)"}
        response = requests.get(url, timeout=15, headers=headers)
        html_content = response.text

        title = ""
        title_start = html_content.find("<title>")
        title_end = html_content.find("</title>")
        if title_start != -1 and title_end != -1:
            title = html_content[title_start + 7: title_end].strip()

        description = ""
        desc_idx = html_content.lower().find('name="description"')
        if desc_idx == -1:
            desc_idx = html_content.lower().find("name='description'")
        if desc_idx != -1:
            content_start = html_content.lower().find('content="', desc_idx)
            if content_start != -1:
                content_end = html_content.find('"', content_start + 9)
                if content_end != -1:
                    description = html_content[content_start + 9: content_end]

        return {
            "title": title,
            "description": description,
            "html_snippet": html_content[:3000],
            "status_code": response.status_code
        }
    except Exception as e:
        return {"error": str(e), "title": "Web App Portal", "description": ""}

# ── AI URL Analyzer ──────────────────────────────────────────────────────────
def analyze_url_ai(url_metadata: dict, url: str) -> dict:
    """Runs a targeted prompt to analyze a live site's purpose/tech stack."""
    if ai_engine.engine is None:
        raise Exception("AI Engine is currently unavailable")

    prompt = f"""You are reviewing a live web application.
URL: {url}

METADATA:
Title: {url_metadata.get('title', '')}
Description: {url_metadata.get('description', '')}

HTML STRUCTURE SNIPPET:
{url_metadata.get('html_snippet', '')}

Respond in strict JSON with these fields:
- "project_type": short label (e.g. "React Web Application", "Streamlit dashboard")
- "purpose": one paragraph detailing what this application does
- "tech_stack": list of estimated technologies used (e.g. React, Tailwind CSS, Python)
- "entry_point": "URL Link"
- "run_command": "Visit {url}"
- "deploy_target": "other"
- "missing_pieces": list of anything that could be improved or is missing
"""
    raw = ai_engine.engine.ask(prompt, json_mode=True)
    return json.loads(raw)

# ── Multi-File Content Extractor ──────────────────────────────────────────────
def extract_text_from_file(uploaded_file) -> str:
    """Extract text content from various file types."""
    name = uploaded_file.name.lower()
    try:
        if name.endswith((".txt", ".py", ".js", ".ts", ".html", ".css",
                           ".json", ".md", ".yaml", ".yml", ".toml", ".ini",
                           ".env", ".sh", ".bat", ".xml", ".csv")):
            return uploaded_file.read().decode("utf-8", errors="ignore")

        elif name.endswith(".pdf"):
            try:
                import PyPDF2
                reader = PyPDF2.PdfReader(io.BytesIO(uploaded_file.read()))
                text = ""
                for page in reader.pages:
                    text += page.extract_text() or ""
                return text[:8000]
            except ImportError:
                return f"[PDF file: {uploaded_file.name} — PyPDF2 not available]"

        elif name.endswith(".docx"):
            try:
                from docx import Document
                doc = Document(io.BytesIO(uploaded_file.read()))
                return "\n".join(para.text for para in doc.paragraphs)[:8000]
            except ImportError:
                return f"[DOCX file: {uploaded_file.name} — python-docx not available]"

        elif name.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(uploaded_file.read()))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text[:8000]
            except ImportError:
                return f"[PPTX file: {uploaded_file.name} — python-pptx not available]"

        elif name.endswith((".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg")):
            return f"[Image file: {uploaded_file.name}]"

        else:
            # Try reading as text anyway
            try:
                return uploaded_file.read().decode("utf-8", errors="ignore")[:5000]
            except Exception:
                return f"[Binary file: {uploaded_file.name}]"
    except Exception as e:
        return f"[Could not read {uploaded_file.name}: {e}]"

# ── Build Synthetic Project from Uploaded Files ───────────────────────────────
def build_project_from_files(uploaded_files, work_dir: str) -> tuple:
    """
    Creates a synthetic file tree and key_files dict from uploaded files.
    Also saves files to work_dir for GitHub push.
    Returns (file_tree, key_files)
    """
    file_tree = []
    key_files = {}
    KEY_NAMES = {
        "requirements.txt", "package.json", "app.py", "main.py",
        "streamlit_app.py", "index.js", "index.html", "README.md",
        "Procfile", "vercel.json", "Dockerfile",
    }

    for uf in uploaded_files:
        filename = uf.name
        file_tree.append(filename)

        # Save to disk for GitHub push
        dest = os.path.join(work_dir, filename)
        uf.seek(0)
        with open(dest, "wb") as f:
            f.write(uf.read())

        # Extract text for AI
        uf.seek(0)
        if filename in KEY_NAMES or any(
            filename.endswith(ext) for ext in
            [".txt", ".py", ".js", ".ts", ".html", ".css", ".json",
             ".md", ".yaml", ".yml", ".toml", ".pdf", ".docx", ".pptx"]
        ):
            text = extract_text_from_file(uf)
            if text:
                key_files[filename] = text[:3000]

    return file_tree, key_files

# ── Premium Theme CSS ─────────────────────────────────────────────────────────
st.markdown(
    """
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;500;600;700&display=swap');

    html, body, [class*="css"] {
        font-family: 'Outfit', sans-serif;
    }

    /* Dark sidebar styling */
    section[data-testid="stSidebar"] {
        background: linear-gradient(180deg, #090d16 0%, #111827 100%) !important;
        border-right: 1px solid #1f2937;
    }
    section[data-testid="stSidebar"] * {
        color: #e5e7eb !important;
    }

    /* Custom main app body background */
    .stApp {
        background-color: #0b0f19;
    }

    /* Buttons */
    .stButton > button {
        background: linear-gradient(135deg, #6366f1 0%, #4f46e5 100%);
        color: white !important;
        border: none;
        border-radius: 8px;
        font-weight: 600;
        padding: 0.6rem 1.8rem;
        transition: all 0.2s ease;
    }
    .stButton > button:hover {
        background: linear-gradient(135deg, #818cf8 0%, #6366f1 100%);
        transform: translateY(-1px);
        box-shadow: 0 4px 15px rgba(99, 102, 241, 0.4);
    }

    /* Cards */
    .metric-card {
        background: linear-gradient(135deg, #111827 0%, #0f172a 100%);
        border: 1px solid #1f2937;
        border-radius: 12px;
        padding: 1.5rem;
        color: #f3f4f6;
    }

    /* Input method panel */
    .input-panel {
        background: linear-gradient(135deg, #111827 0%, #0f172a 100%);
        border: 1px solid #1e293b;
        border-radius: 12px;
        padding: 1.5rem;
        margin-top: 0.5rem;
    }

    h1, h2, h3, h4, h5, h6 {
        color: #f3f4f6 !important;
    }

    /* Highlighted link tabs */
    a.link-btn {
        display: inline-block;
        background: linear-gradient(135deg, #10b981 0%, #059669 100%);
        color: white !important;
        text-decoration: none;
        padding: 0.5rem 1.2rem;
        border-radius: 6px;
        font-weight: 600;
        margin-right: 10px;
        transition: all 0.2s;
    }
    a.link-btn:hover {
        background: linear-gradient(135deg, #34d399 0%, #10b981 100%);
        box-shadow: 0 4px 12px rgba(16, 185, 129, 0.3);
    }

    /* Progress step styling */
    .step-badge {
        display: inline-block;
        background: linear-gradient(135deg, #6366f1, #4f46e5);
        color: white;
        border-radius: 20px;
        padding: 2px 10px;
        font-size: 0.78rem;
        font-weight: 600;
        margin-right: 6px;
    }

    /* Radio button improvements */
    .stRadio > div {
        gap: 0.5rem;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

# ── Sidebar Configurations ────────────────────────────────────────────────────
st.sidebar.markdown(
    """
    <div style="text-align:center;padding:1rem 0;">
        <div style="font-size:2.5rem">🚀</div>
        <div style="font-weight:700;font-size:1.3rem;color:#f3f4f6;letter-spacing:1px;">
            DevLaunch AI
        </div>
        <div style="font-size:0.8rem;color:#9ca3af;margin-top:0.2rem;">Automation Agent</div>
    </div>
    <hr style="border-color:#1f2937;margin:0.5rem 0 1.5rem;">
    """,
    unsafe_allow_html=True,
)

# Load env/secrets defaults
def _secret(key: str) -> str:
    try:
        return st.secrets.get(key, "") or os.getenv(key, "")
    except Exception:
        return os.getenv(key, "")

env_groq = _secret("GROQ_API_KEY")
env_github = _secret("GITHUB_TOKEN")
env_vercel = _secret("VERCEL_TOKEN")
env_username = _secret("GITHUB_USERNAME")

st.sidebar.subheader("🔑 API Configurations")
groq_key = st.sidebar.text_input(
    "Groq API Key", value=env_groq, type="password",
    help="Used to analyze project and generate PRD/Marketing copy"
)
github_token = st.sidebar.text_input(
    "GitHub Token", value=env_github, type="password",
    help="Used to create repository and push files"
)
github_username = st.sidebar.text_input(
    "GitHub Username", value=env_username,
    placeholder="e.g. shabihaiqbal0"
)
vercel_token = st.sidebar.text_input(
    "Vercel Token (Optional)", value=env_vercel, type="password",
    help="Used to trigger deployments"
)

# Apply dynamic reinitialization
reinit_clients(groq_key, github_token, vercel_token)

# ── Hero Header ───────────────────────────────────────────────────────────────
st.markdown(
    """
    <div style="
        background: linear-gradient(135deg, #1e1b4b 0%, #0b0f19 50%, #030712 100%);
        border: 1px solid #1e293b;
        border-radius: 16px;
        padding: 2.5rem 2rem;
        margin-bottom: 2rem;
        box-shadow: 0 4px 20px rgba(0, 0, 0, 0.4);
    ">
        <h1 style="font-size:2.8rem; margin:0; font-weight: 700; background: linear-gradient(to right, #a5b4fc, #818cf8); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">
            🚀 DevLaunch AI Workspace Agent
        </h1>
        <p style="color:#9ca3af; font-size:1.15rem; margin-top:0.6rem; max-width:800px;">
            A complete automation agent: choose any project input to analyze, generate high-fidelity documentation, capture breakpoint screenshots, and push to GitHub or deploy automatically.
        </p>
    </div>
    """,
    unsafe_allow_html=True,
)

# ── Main Flow Form ────────────────────────────────────────────────────────────
col_left, col_right = st.columns([2, 3])

with col_left:
    st.subheader("📁 Choose Input Method")

    input_method = st.radio(
        "Select how you want to provide your project:",
        ["📄 Upload Files", "📦 Upload Folder (ZIP)", "🌐 Enter URL"],
        help="Choose the method that's easiest for you."
    )

    # ── Panel: Upload Files ──────────────────────────────────────────────────
    if input_method == "📄 Upload Files":
        st.markdown(
            '<div class="input-panel">',
            unsafe_allow_html=True,
        )
        st.markdown("**Upload one or multiple project files:**")
        uploaded_files = st.file_uploader(
            "Supported: Images, PDF, DOCX, PPTX, TXT, PY, JS, HTML, CSS, JSON, ZIP, and more",
            type=[
                "txt", "py", "js", "ts", "jsx", "tsx",
                "html", "css", "json", "yaml", "yml", "toml", "md",
                "pdf", "docx", "pptx",
                "png", "jpg", "jpeg", "gif", "webp",
                "zip", "sh", "bat", "xml", "csv", "ini", "env"
            ],
            accept_multiple_files=True,
            key="multi_file_uploader",
        )
        if uploaded_files:
            st.success(f"✅ {len(uploaded_files)} file(s) selected")
            for f in uploaded_files:
                st.caption(f"📎 {f.name} ({f.size:,} bytes)")
        st.markdown("</div>", unsafe_allow_html=True)

    # ── Panel: Upload Folder (ZIP) ───────────────────────────────────────────
    elif input_method == "📦 Upload Folder (ZIP)":
        st.markdown('<div class="input-panel">', unsafe_allow_html=True)
        st.markdown("**Zip your entire project folder and upload it:**")
        st.caption("💡 Right-click your folder → Send to → Compressed (ZIP) folder")
        zip_upload = st.file_uploader(
            "Upload Project ZIP",
            type=["zip"],
            key="folder_zip_uploader",
        )
        if zip_upload:
            st.success(f"✅ ZIP uploaded: {zip_upload.name} ({zip_upload.size:,} bytes)")
        st.markdown("</div>", unsafe_allow_html=True)

    # ── Panel: Enter URL ─────────────────────────────────────────────────────
    elif input_method == "🌐 Enter URL":
        st.markdown('<div class="input-panel">', unsafe_allow_html=True)
        st.markdown("**Enter any public URL:**")
        url_input = st.text_input(
            "URL",
            placeholder="e.g. https://myapp.streamlit.app, https://myapp.vercel.app",
            key="url_input_field",
        )
        st.caption("Works with: Streamlit apps, Vercel, Netlify, GitHub Pages, any public web app")
        if url_input:
            st.info(f"🔗 Target: {url_input}")
        st.markdown("</div>", unsafe_allow_html=True)

    st.markdown("---")
    st.subheader("⚙️ Deployment Settings")
    repo_name = st.text_input("Repository Name", placeholder="e.g. my-cool-project")
    make_private = st.checkbox("Make GitHub Repository Private", value=False)
    include_content = st.checkbox("Generate Marketing & Publishing Content", value=True)
    capture_screenshot = st.checkbox("Capture App Screenshots (Playwright)", value=True)

    run_pipeline_btn = st.button("🚀 Launch Automation Pipeline", use_container_width=True)

with col_right:
    st.subheader("📋 Execution & Live Console")
    console_placeholder = st.empty()
    progress_placeholder = st.empty()
    console_placeholder.info("Ready. Set configuration and click 'Launch Automation Pipeline' to begin.")

    if run_pipeline_btn:
        # ── Validation ───────────────────────────────────────────────────────
        has_errors = False

        if not groq_key:
            st.error("❌ Groq API Key is required.")
            has_errors = True

        is_url_source = input_method == "🌐 Enter URL"
        requires_github = not is_url_source and bool(github_token and github_username and repo_name)

        # Specific source validations
        if input_method == "📄 Upload Files":
            if not uploaded_files:
                st.error("❌ Please upload at least one file.")
                has_errors = True
        elif input_method == "📦 Upload Folder (ZIP)":
            if not zip_upload:
                st.error("❌ Please upload a ZIP file.")
                has_errors = True
        elif input_method == "🌐 Enter URL":
            if not url_input:
                st.error("❌ Please enter a URL.")
                has_errors = True
            elif not (url_input.startswith("http://") or url_input.startswith("https://")):
                st.error("❌ URL must start with http:// or https://")
                has_errors = True

        if not has_errors:
            console_log = []

            def add_log(msg: str, status: str = "info"):
                console_log.append(f"[{datetime.now().strftime('%H:%M:%S')}] {msg}")
                log_text = "\n".join(console_log)
                if status == "success":
                    console_placeholder.code(log_text + "\n\n✅ Pipeline completed successfully!")
                else:
                    console_placeholder.code(log_text)

            add_log("Starting Automation Pipeline...")

            # Setup temp workspace
            temp_dir_obj = tempfile.TemporaryDirectory()
            extract_path = os.path.join(temp_dir_obj.name, "extracted")
            os.makedirs(extract_path, exist_ok=True)
            work_dir = extract_path

            target_url = url_input if is_url_source else ""

            try:
                # ── Step 0: Load / Prepare Source ──
                progress_placeholder.progress(5, text="📂 Preparing source...")

                if input_method == "📄 Upload Files":
                    add_log(f"Processing {len(uploaded_files)} uploaded file(s)...")
                    file_tree, key_files = build_project_from_files(uploaded_files, work_dir)
                    add_log(f"Extracted content from: {', '.join(f.name for f in uploaded_files[:5])}")

                    # Check if any ZIP was uploaded — extract it
                    for uf in uploaded_files:
                        if uf.name.lower().endswith(".zip"):
                            uf.seek(0)
                            zip_extract_dir = os.path.join(work_dir, "zip_contents")
                            os.makedirs(zip_extract_dir, exist_ok=True)
                            with zipfile.ZipFile(io.BytesIO(uf.read()), 'r') as zf:
                                zf.extractall(zip_extract_dir)
                            add_log(f"ZIP '{uf.name}' also extracted to workspace.")

                elif input_method == "📦 Upload Folder (ZIP)":
                    add_log("Extracting uploaded ZIP archive...")
                    progress_placeholder.progress(10, text="📦 Extracting ZIP...")
                    with zipfile.ZipFile(io.BytesIO(zip_upload.read()), 'r') as zf:
                        zf.extractall(extract_path)

                    # If ZIP contains a single root folder, descend into it
                    contents = os.listdir(extract_path)
                    if len(contents) == 1 and os.path.isdir(os.path.join(extract_path, contents[0])):
                        work_dir = os.path.join(extract_path, contents[0])
                    add_log(f"ZIP extracted. Workspace ready: {len(contents)} item(s) at root.")

                elif input_method == "🌐 Enter URL":
                    add_log(f"Fetching metadata from URL: {target_url}...")

                # ── Step 1: Scan & Analyze ──
                progress_placeholder.progress(20, text="🔍 Analyzing project...")
                add_log("Step 1: Analyzing project with AI...")

                if is_url_source:
                    add_log(f"Scraping HTML metadata from {target_url}...")
                    url_meta = scrape_url_metadata(target_url)
                    if "error" in url_meta and not url_meta.get("title"):
                        raise Exception(f"Could not reach URL: {url_meta['error']}")
                    add_log(f"Page title detected: '{url_meta.get('title', 'N/A')}'")
                    add_log("Running AI analysis on page structure...")
                    analysis_res = analyze_url_ai(url_meta, target_url)
                    analysis = {
                        "file_tree": ["index.html"],
                        "key_files_found": ["index.html"],
                        "ai_summary": analysis_res
                    }

                elif input_method == "📄 Upload Files":
                    # Run AI directly on extracted content
                    if ai_engine.engine is None:
                        raise Exception("AI Engine unavailable — check Groq API key")
                    add_log("Running AI analysis on uploaded file contents...")
                    ai_summary = ai_engine.engine.analyze_project(
                        file_tree=file_tree,
                        key_files=key_files
                    )
                    analysis = {
                        "file_tree": file_tree,
                        "key_files_found": list(key_files.keys()),
                        "ai_summary": ai_summary
                    }

                else:  # ZIP folder
                    analysis = analyze(work_dir)
                    add_log("Framework scanning complete.")

                add_log(f"✓ Project type: {analysis['ai_summary'].get('project_type', 'Unknown')}")
                add_log(f"✓ Tech stack: {', '.join(analysis['ai_summary'].get('tech_stack', []))}")

                # ── Step 2: Push to GitHub ──
                progress_placeholder.progress(40, text="📤 Pushing to GitHub...")
                github_url = None
                if requires_github:
                    add_log("Step 2: Authenticating and pushing to GitHub...")
                    try:
                        client = github_client_mod.github_client or github_client_mod.GitHubClient(token=github_token)
                        github_url = client.push_project(
                            project_path=work_dir,
                            repo_name=repo_name,
                            username=github_username,
                            private=make_private,
                        )
                        add_log(f"✓ GitHub repository: {github_url}")
                    except Exception as e:
                        add_log(f"⚠️ GitHub push failed/skipped: {e}")
                else:
                    add_log("Step 2: GitHub push skipped (no credentials or URL source).")

                # ── Step 3: Deploy ──
                progress_placeholder.progress(55, text="🚀 Deploying...")
                live_url = None
                if is_url_source:
                    live_url = target_url
                    add_log(f"Step 3: App is already live at: {live_url}")
                elif vercel_token and vercel_client_mod.vercel_client is not None:
                    add_log("Step 3: Deploying to Vercel...")
                    try:
                        live_url = vercel_client_mod.vercel_client.deploy(
                            repo_name=repo_name,
                            github_username=github_username,
                        )
                        add_log(f"✓ Vercel deployment: {live_url}")
                    except Exception as e:
                        add_log(f"⚠️ Vercel deployment skipped: {e}")
                else:
                    add_log("Step 3: Deployment skipped (Vercel token not configured).")

                # ── Step 4: Generate Documentation ──
                progress_placeholder.progress(70, text="📝 Generating docs...")
                add_log("Step 4: Compiling PRD, README, and social copies...")
                summary = build_summary(
                    project_name=repo_name if repo_name else "web-app",
                    analysis=analysis["ai_summary"],
                    github_url=github_url,
                    live_url=live_url,
                    include_content=include_content,
                )

                social_drafts = {}
                if include_content:
                    try:
                        import backend.screenshot.formatter as formatter_mod
                        if formatter_mod.formatter is not None:
                            social_drafts = formatter_mod.formatter.format_all(
                                analysis["ai_summary"], live_url=live_url
                            )
                    except Exception as e:
                        add_log(f"⚠️ Social formatting error: {e}")
                add_log("✓ Documentation compiler finished.")

                # ── Step 5: Screenshots ──
                progress_placeholder.progress(85, text="📸 Capturing screenshots...")
                screenshots = []
                capture_target = live_url if live_url else target_url
                if capture_screenshot and capture_target:
                    add_log(f"Step 5: Capturing screenshots of '{capture_target}'...")
                    try:
                        screenshot_dir = os.path.join(temp_dir_obj.name, "screenshots")
                        os.makedirs(screenshot_dir, exist_ok=True)
                        screenshots = capture_screenshots(capture_target, screenshot_dir)
                        if screenshots:
                            add_log(f"✓ {len(screenshots)} screenshot(s) captured.")
                        else:
                            add_log("⚠️ Screenshots skipped (Playwright/Chromium not available in this environment).")
                    except Exception as e:
                        add_log(f"⚠️ Screenshot capture failed gracefully: {e}")
                else:
                    add_log("Step 5: Screenshot capture skipped.")

                # ── Done ──
                progress_placeholder.progress(100, text="✅ Pipeline complete!")
                add_log("Pipeline completed!", status="success")

                # Copy screenshots to a persistent location in session
                persistent_screenshots = []
                for s in screenshots:
                    try:
                        dest_name = os.path.basename(s)
                        dest_path = os.path.join(ROOT, "storage", "screenshots", dest_name)
                        os.makedirs(os.path.dirname(dest_path), exist_ok=True)
                        shutil.copy2(s, dest_path)
                        persistent_screenshots.append(dest_path)
                    except Exception:
                        persistent_screenshots.append(s)

                st.session_state["pipeline_results"] = {
                    "summary": summary.to_dict(),
                    "github_url": github_url,
                    "live_url": live_url,
                    "screenshots": persistent_screenshots,
                    "project_name": repo_name if repo_name else "Web Portal",
                    "social_drafts": social_drafts,
                }
                st.success("🎉 Pipeline execution successful!")
                st.balloons()

            except Exception as e:
                add_log(f"❌ Execution failed: {e}")
                st.error(f"Pipeline error: {e}")
                progress_placeholder.empty()
            finally:
                try:
                    temp_dir_obj.cleanup()
                except Exception:
                    pass

# ── Results Sections ──────────────────────────────────────────────────────────
if "pipeline_results" in st.session_state:
    st.markdown("---")
    res = st.session_state["pipeline_results"]
    summary = res["summary"]

    st.subheader(f"📊 Results for {res['project_name']}")

    # URL Actions
    col_urls = st.columns(3)
    with col_urls[0]:
        if res["github_url"]:
            st.markdown(
                f'<a href="{res["github_url"]}" target="_blank" class="link-btn">📁 Open GitHub Repository</a>',
                unsafe_allow_html=True
            )
        else:
            st.info("GitHub Repository was not created.")
    with col_urls[1]:
        if res["live_url"]:
            st.markdown(
                f'<a href="{res["live_url"]}" target="_blank" class="link-btn" style="background:linear-gradient(135deg, #6366f1, #4f46e5)">🌐 View Live Deployment</a>',
                unsafe_allow_html=True
            )
        else:
            st.info("No live deployment URL available.")

    # Tabs
    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "📝 Product Requirements (PRD)",
        "📖 README.md",
        "📢 Marketing & Social Content",
        "📂 Tech Stack Analysis",
        "📸 Screenshots"
    ])

    with tab1:
        st.subheader("Product Requirements Document (PRD)")
        if summary.get("prd"):
            st.download_button("⬇️ Download PRD.md", summary["prd"], file_name="PRD.md", mime="text/markdown")
            st.markdown(summary["prd"])
        else:
            st.info("PRD was not generated.")

    with tab2:
        st.subheader("README.md File Content")
        if summary.get("readme"):
            st.download_button("⬇️ Download README.md", summary["readme"], file_name="README.md", mime="text/markdown")
            st.code(summary["readme"], language="markdown")
        else:
            st.info("README content not generated.")

    with tab3:
        st.subheader("Social Drafts & Gig Descriptions")
        social_drafts = res.get("social_drafts", {})

        col_soc_l, col_soc_r = st.columns(2)
        with col_soc_l:
            st.markdown("#### 👥 LinkedIn Post")
            linkedin_post = social_drafts.get("linkedin") or "N/A"
            st.text_area("LinkedIn Draft", linkedin_post, height=180)

            st.markdown("#### 💳 Fiverr Gig Portfolio Blurb")
            fiverr_gig = social_drafts.get("fiverr") or "N/A"
            st.text_area("Fiverr Draft", fiverr_gig, height=120)

        with col_soc_r:
            st.markdown("#### 💼 Upwork Portfolio Description")
            upwork_port = social_drafts.get("upwork") or "N/A"
            st.text_area("Upwork Draft", upwork_port, height=180)

            st.markdown("#### ✍️ Dev Blog Post")
            blog_post = summary.get("blog_post") or "N/A"
            st.text_area("Blog Post Draft", blog_post, height=180)

    with tab4:
        st.subheader("Project Technical Profile")
        ai_sum = summary.get("analysis", {})
        if ai_sum:
            col_tech1, col_tech2 = st.columns(2)
            with col_tech1:
                st.write(f"**Project Type**: {ai_sum.get('project_type', 'N/A')}")
                st.write(f"**Main Entry Point**: `{ai_sum.get('entry_point', 'N/A')}`")
                st.write(f"**Run Command**: `{ai_sum.get('run_command', 'N/A')}`")
                st.write(f"**Deploy Target**: `{ai_sum.get('deploy_target', 'N/A')}`")
            with col_tech2:
                st.write("**Tech Stack:**")
                for tech in ai_sum.get("tech_stack", []):
                    st.markdown(f"- {tech}")

            st.write("---")
            st.markdown("#### Purpose")
            st.write(ai_sum.get("purpose", "N/A"))

            st.markdown("#### Missing Pieces")
            missing = ai_sum.get("missing_pieces", [])
            if missing:
                for item in missing:
                    st.markdown(f"⚠️ {item}")
            else:
                st.success("No missing configuration elements detected! Ready to compile and ship.")
        else:
            st.info("Tech stack profile is not available.")

    with tab5:
        st.subheader("App Layout Breakpoint Captures")
        if res.get("screenshots"):
            for s_path in res["screenshots"]:
                if os.path.isfile(s_path):
                    name = os.path.basename(s_path).replace("screenshot_", "").replace(".png", "").capitalize()
                    st.image(s_path, caption=f"{name} Breakpoint Preview", use_container_width=True)
        else:
            st.info("No screenshots were captured. Screenshots require Playwright/Chromium and a live URL.")
